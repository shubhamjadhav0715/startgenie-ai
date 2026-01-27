"""
PowerPoint Generator
Creates professional PPT presentations from blueprints
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from typing import Dict, Any


class PPTGenerator:
    """Generate PowerPoint presentations from blueprint data"""
    
    def __init__(self):
        self.primary_color = RGBColor(2, 132, 199)  # #0284c7
        self.secondary_color = RGBColor(3, 105, 161)  # #0369a1
        self.text_color = RGBColor(31, 41, 55)  # #1f2937
    
    def generate_blueprint_ppt(self, blueprint_data: Dict[str, Any]) -> BytesIO:
        """
        Generate PowerPoint from blueprint data
        
        Args:
            blueprint_data: Blueprint dictionary with content
        
        Returns:
            BytesIO buffer containing PPT
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        content = blueprint_data.get('content', {})
        
        # Title slide
        self._create_title_slide(prs, blueprint_data)
        
        # Overview
        if content.get('startup_overview'):
            self._create_overview_slides(prs, content['startup_overview'])
        
        # Market Analysis
        if content.get('market_analysis'):
            self._create_market_slides(prs, content['market_analysis'])
        
        # Business Model
        if content.get('business_model'):
            self._create_business_model_slide(prs, content['business_model'])
        
        # SWOT
        if content.get('swot_analysis'):
            self._create_swot_slide(prs, content['swot_analysis'])
        
        # Budget
        if content.get('budget_estimation'):
            self._create_budget_slide(prs, content['budget_estimation'])
        
        # Funding
        if content.get('funding_investment'):
            self._create_funding_slide(prs, content['funding_investment'])
        
        # Roadmap
        if content.get('action_roadmap'):
            self._create_roadmap_slide(prs, content['action_roadmap'])
        
        # Thank you slide
        self._create_thank_you_slide(prs)
        
        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    def _create_title_slide(self, prs: Presentation, blueprint_data: Dict[str, Any]):
        """Create title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Startup Blueprint"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.primary_color
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = blueprint_data.get('startup_idea', '')
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Powered by StartGenie AI"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.color.rgb = RGBColor(107, 114, 128)
    
    def _create_overview_slides(self, prs: Presentation, data: Dict[str, Any]):
        """Create overview slides"""
        # Problem & Solution slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        title = slide.shapes.title
        title.text = "Problem & Solution"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        # Problem
        p = content.add_paragraph()
        p.text = "Problem"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.space_after = Pt(10)
        
        p = content.add_paragraph()
        p.text = data.get('problem_statement', 'N/A')
        p.font.size = Pt(16)
        p.space_after = Pt(20)
        
        # Solution
        p = content.add_paragraph()
        p.text = "Solution"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.space_after = Pt(10)
        
        p = content.add_paragraph()
        p.text = data.get('solution', 'N/A')
        p.font.size = Pt(16)
        
        # UVP slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Unique Value Proposition"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        p = content.add_paragraph()
        p.text = data.get('unique_value_proposition', 'N/A')
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.CENTER
    
    def _create_market_slides(self, prs: Presentation, data: Dict[str, Any]):
        """Create market analysis slides"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Market Analysis"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        # Target Audience
        p = content.add_paragraph()
        p.text = "Target Audience"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        
        p = content.add_paragraph()
        p.text = data.get('target_audience', 'N/A')
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(15)
        
        # Market Size
        p = content.add_paragraph()
        p.text = "Market Size"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        
        p = content.add_paragraph()
        p.text = data.get('market_size', 'N/A')
        p.font.size = Pt(14)
        p.level = 1
        
        # Industry Trends slide
        if data.get('industry_trends'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Industry Trends"
            
            content = slide.placeholders[1].text_frame
            content.clear()
            
            for trend in data['industry_trends']:
                p = content.add_paragraph()
                p.text = trend
                p.font.size = Pt(16)
                p.level = 1
                p.space_after = Pt(12)
    
    def _create_business_model_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Create business model slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Business Model"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        # Revenue Streams
        p = content.add_paragraph()
        p.text = "Revenue Streams"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.space_after = Pt(10)
        
        for stream in data.get('revenue_streams', []):
            p = content.add_paragraph()
            p.text = stream
            p.font.size = Pt(14)
            p.level = 1
            p.space_after = Pt(8)
        
        # Pricing Strategy
        p = content.add_paragraph()
        p.text = "Pricing Strategy"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.space_before = Pt(15)
        p.space_after = Pt(10)
        
        p = content.add_paragraph()
        p.text = data.get('pricing_strategy', 'N/A')
        p.font.size = Pt(14)
        p.level = 1
    
    def _create_swot_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Create SWOT analysis slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "SWOT Analysis"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.primary_color
        
        # Create 4 quadrants
        quadrants = [
            ("Strengths", data.get('strengths', []), Inches(0.5), Inches(1.5), RGBColor(34, 197, 94)),
            ("Weaknesses", data.get('weaknesses', []), Inches(5.25), Inches(1.5), RGBColor(239, 68, 68)),
            ("Opportunities", data.get('opportunities', []), Inches(0.5), Inches(4.5), RGBColor(59, 130, 246)),
            ("Threats", data.get('threats', []), Inches(5.25), Inches(4.5), RGBColor(251, 191, 36))
        ]
        
        for title, items, left, top, color in quadrants:
            box = slide.shapes.add_textbox(left, top, Inches(4.25), Inches(2.5))
            frame = box.text_frame
            frame.word_wrap = True
            
            # Title
            p = frame.add_paragraph()
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = color
            p.space_after = Pt(10)
            
            # Items
            for item in items[:4]:  # Limit to 4 items per quadrant
                p = frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.level = 1
                p.space_after = Pt(6)
    
    def _create_budget_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Create budget slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Budget Estimation"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        budget_items = [
            ("Initial Setup Cost", data.get('initial_setup_cost')),
            ("Monthly Operational Expenses", data.get('monthly_operational_expenses')),
            ("Technology Cost", data.get('technology_cost')),
            ("Marketing Cost", data.get('marketing_cost'))
        ]
        
        for label, value in budget_items:
            if value:
                p = content.add_paragraph()
                p.text = f"{label}: ₹{value:,.0f}"
                p.font.size = Pt(18)
                p.level = 1
                p.space_after = Pt(15)
    
    def _create_funding_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Create funding slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Funding Options"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        for option in data.get('funding_options', []):
            p = content.add_paragraph()
            p.text = option
            p.font.size = Pt(16)
            p.level = 1
            p.space_after = Pt(12)
    
    def _create_roadmap_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Create roadmap slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Action Roadmap"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        phases = [
            ("Months 0-3", data.get('months_0_3', [])),
            ("Months 3-6", data.get('months_3_6', [])),
            ("Months 6-12", data.get('months_6_12', []))
        ]
        
        for phase_name, actions in phases:
            p = content.add_paragraph()
            p.text = phase_name
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            p.space_after = Pt(8)
            
            for action in actions[:2]:  # Limit to 2 actions per phase
                p = content.add_paragraph()
                p.text = action
                p.font.size = Pt(12)
                p.level = 1
                p.space_after = Pt(6)
            
            p.space_after = Pt(15)
    
    def _create_thank_you_slide(self, prs: Presentation):
        """Create thank you slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Thank you text
        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        text_frame = text_box.text_frame
        text_frame.text = "Thank You"
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(54)
        text_para.font.bold = True
        text_para.font.color.rgb = self.primary_color
        text_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Generated by StartGenie AI"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(18)
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.color.rgb = RGBColor(107, 114, 128)
